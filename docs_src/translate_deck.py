# -*- coding: utf-8 -*-
"""Produce the English deck -> public/presentation_2503_en.pptx

Translates public/presentation_2503.pptx in place from deck_strings.T, so the
original design, layout and imagery are preserved exactly. Text is replaced at
paragraph level (the first run keeps its formatting and receives the whole
translated string; sibling runs are emptied), because a sentence is often split
across several runs by stray formatting.
"""
import os
import re
import shutil
import subprocess
import sys

from pptx import Presentation

from deck_strings import T

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(BASE, "public", "presentation_2503.pptx")
OUT = os.path.join(BASE, "public", "presentation_2503_en.pptx")

CYRILLIC = re.compile(r"[Ѐ-ӿ]")


def iter_paragraphs(shape):
    """Every paragraph in a shape, descending into tables and groups."""
    if shape.shape_type == 6 and hasattr(shape, "shapes"):        # group
        for inner in shape.shapes:
            yield from iter_paragraphs(inner)
        return
    if shape.has_text_frame:
        yield from shape.text_frame.paragraphs
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame.paragraphs


def translate(path_in, path_out):
    prs = Presentation(path_in)
    replaced = 0
    missing = []

    for index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            for para in iter_paragraphs(shape):
                runs = para.runs
                if not runs:
                    continue
                whole = "".join(r.text for r in runs)
                key = whole.strip()
                if not key:
                    continue
                if key in T:
                    # keep leading/trailing spacing exactly as authored
                    lead = whole[: len(whole) - len(whole.lstrip())]
                    tail = whole[len(whole.rstrip()):]
                    runs[0].text = lead + T[key] + tail
                    for r in runs[1:]:
                        r.text = ""
                    replaced += 1
                elif CYRILLIC.search(key):
                    missing.append((index, key))

    prs.save(path_out)
    print(f"replaced {replaced} paragraphs -> {path_out}")
    if missing:
        print(f"\n{len(missing)} untranslated Russian paragraphs remain:")
        for slide_no, text in missing:
            print(f"  slide {slide_no}: {text[:90]!r}")
        return 1
    print("no Russian text left in the deck")
    return 0


SOFFICE_CANDIDATES = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "soffice",
    "libreoffice",
]


def export_pdf(pptx_path):
    """Render the deck to PDF with LibreOffice, if it is available.

    PowerPoint/Keynote can do this too, but only through GUI automation, which
    needs macOS Automation and Accessibility grants. LibreOffice converts
    headlessly with no such permissions.
    """
    exe = next((c for c in SOFFICE_CANDIDATES
                if os.path.exists(c) or shutil.which(c)), None)
    if exe is None:
        print("LibreOffice not found - skipping PDF export.\n"
              "  brew install --cask libreoffice   (or export from PowerPoint)")
        return False
    outdir = os.path.dirname(pptx_path)
    subprocess.run([exe, "--headless", "--norestore", "--convert-to", "pdf",
                    "--outdir", outdir, pptx_path],
                   check=True, capture_output=True)
    pdf = os.path.splitext(pptx_path)[0] + ".pdf"
    print(f"written: {pdf}")
    return True


if __name__ == "__main__":
    status = translate(SRC, OUT)
    export_pdf(OUT)
    sys.exit(status)
